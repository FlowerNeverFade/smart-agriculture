# -*- coding: utf-8 -*-
"""Rebuild part 3 (技术架构与数据主线) of the defense deck into a copy.

Source template/assets: 桌面 PPT郑.pptx. Output: 桌面 PPT郑_第三部分重制.pptx
Operations: reorder part-3 slides, add 2 new slides (数据主线九跳 / 工程可信度),
add bottom message strips on diagram slides, move the dense 6-layer map to
appendix, and write speaker notes for slides 16-22 + appendix.
"""
import copy
import shutil

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SRC = r"C:\Users\Yang\Desktop\PPT原.pptx"  # 2026-09-03 由 PPT郑.pptx 改名
DST = r"C:\Users\Yang\Desktop\PPT郑_第三部分重制.pptx"

PX = 9525
BLUE = RGBColor(0x1F, 0x4E, 0x79)
BLUE2 = RGBColor(0x1F, 0x6F, 0xB2)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
ORANGE = RGBColor(0xE8, 0x77, 0x22)
INK = RGBColor(0x16, 0x36, 0x5C)
GRAY = RGBColor(0x5A, 0x6B, 0x7B)
LIGHT = RGBColor(0xEA, 0xF1, 0xF8)
ORANGEBG = RGBColor(0xFF, 0xF4, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "微软雅黑"


def px(v):
    return Emu(int(v * PX))


def cjk(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def style_run(run, size, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    cjk(run)


def fill_text(tf, lines, align=PP_ALIGN.LEFT):
    """lines: list of (text, size, bold, color)."""
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = text
        style_run(r, size, bold, color)


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.margin_left = px(4)
    tf.margin_right = px(4)
    tf.margin_top = px(2)
    tf.margin_bottom = px(2)
    tf.vertical_anchor = anchor
    fill_text(tf, lines, align)
    return box


def card(slide, x, y, w, h, title, body, line_color, fill=WHITE, title_size=12, body_size=9):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h))
    sp.adjustments[0] = 0.08
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line_color
    sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_left = px(7)
    tf.margin_right = px(7)
    tf.margin_top = px(5)
    tf.margin_bottom = px(4)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = [(title, title_size, True, line_color)]
    lines += [(t, body_size, False, RGBColor(0x33, 0x3F, 0x4A)) for t in body]
    fill_text(tf, lines)
    return sp


def arrow(slide, x, y, w, h, color, shape=MSO_SHAPE.RIGHT_ARROW, rot=0):
    sp = slide.shapes.add_shape(shape, px(x), px(y), px(w), px(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    if rot:
        sp.rotation = rot
    return sp


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def clone_chrome(prs, donor, new_slide, title):
    """Copy template triangle + title textbox from donor slide, set new title."""
    spTree = new_slide.shapes._spTree
    tri = tit = None
    for sh in donor.shapes:
        if sh.name == "直角三角形 1":
            tri = sh._element
        if sh.name == "TextBox 4":
            tit = sh._element
    for el in (tri, tit):
        if el is not None:
            spTree.insert(2, copy.deepcopy(el))
    for sh in new_slide.shapes:
        if sh.name == "TextBox 4" and sh.has_text_frame:
            tf = sh.text_frame
            p = tf.paragraphs[0]
            runs = p.runs
            if runs:
                runs[0].text = title
                for r in runs[1:]:
                    r._r.getparent().remove(r._r)
            for extra in tf.paragraphs[1:]:
                extra._p.getparent().remove(extra._p)
            return sh
    return None


def biggest_pic(slide):
    best, best_area = None, -1
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            w, h = sh.width or 0, sh.height or 0
            if w * h > best_area:
                best, best_area = sh, w * h
    return best


def refit_pic(slide, top=129, bottom=622):
    pic = biggest_pic(slide)
    if pic is None:
        return
    w, h = pic.width or 0, pic.height or 0
    if not w or not h:
        from PIL import Image
        import io
        iw, ih = Image.open(io.BytesIO(pic.image.blob)).size
        aspect = iw / ih
    else:
        aspect = w / h
    new_h = bottom - top
    new_w = int(new_h * aspect)
    if new_w > 1200:
        new_w = 1200
        new_h = int(new_w / aspect)
    pic.height = px(new_h)
    pic.width = px(new_w)
    pic.top = px(top)
    pic.left = px((1280 - new_w) // 2)


def strip(slide, text, y=632, h=72, fill=LIGHT, line=BLUE, size=11):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(40), px(y), px(1200), px(h))
    sp.adjustments[0] = 0.16
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line
    sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_left = px(12)
    tf.margin_right = px(12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    style_run(r, size, False, INK)
    return sp


def main():
    shutil.copyfile(SRC, DST)
    prs = Presentation(DST)
    slides = prs.slides
    I_DIV, I_OLD26, I_POS, I_MAIN, I_FE, I_BE = 15, 16, 17, 18, 19, 20

    # ---------- speaker notes on existing slides ----------
    set_notes(slides[I_DIV],
              "[10s] 一句话钩子：前面两章讲了我们“做什么”“有什么创新”，这一章回答“凭什么信它”。"
              "三个关键词：模块化单体、事件驱动、三条主线的技术边界。")
    set_notes(slides[I_POS],
              "[60s] 这页只讲三件事。一，定位：不是再堆一个传感器看板，而是把 计划→感知→核验→诊断→处方→执行→验证 "
              "串成同一条可追溯闭环。二，三支柱：左=模拟/真实数据传输，中=智能体决策，右=可视化呈现与操作，箭头单向、职责不越界。"
              "三，最重要是底部两个框：本期边界=模拟数据+虚拟执行器+本地可复现测试；BearPi E53_IA1 真实遥测按来源标记接入。"
              "主动讲边界不是露怯，是工程诚实。")
    set_notes(slides[I_MAIN],
              "[60s] 三条主线各一句话。数据主线：事件经 MQTT 进、Redis Stream 缓冲、校验去重质量评分后落库、SSE 推出。"
              "智能体主线：RAG 检索 + 确定性服务出诊断和处方，模型只组合与解释。可视化主线：REST 查询 + SSE 订阅，页面秒级刷新。"
              "然后停在三条硬边界上：浏览器不直连 MQTT；模拟器不写告警和结论；模型不拼 SQL、MQTT topic 或任意 HTTP。"
              "这三条是“可信”的技术表达。")
    set_notes(slides[I_BE],
              "[60s] 后端一句话：单进程、逻辑模块、可替换适配器。内核是 AgriEngine，外围 CropPackCatalog、ControlledLearningService、"
              "AdminManagementService、FarmGovernanceService、SimulationEngine 职责清晰；数据源与执行器都是接口"
              "（DataSource / Actuator），模拟器与 BearPi 适配器只是两个实现。"
              "若被问“为什么是一个万行类”：15 天选模块化单体，模块边界是逻辑的、靠领域事件解耦，拆服务是 P2——主动说，不等问。")
    set_notes(slides[I_FE],
              "[50s] 前端三句话。一，三角色三入口 farmer / index / sysadmin，共享同一内核：api.js 管 REST 与 SSE 双模式、"
              "live-data.js 管数据规范化、roles.js 管权限。二，浏览器只走 REST 和 SSE，MQTT 凭据从不下发到前端。"
              "三，诚实讲取舍：Vue 3 全局运行时 + 模块化原生 JS，Vite 构建配置在但没上 TS/SFC，是 15 天的取舍不是疏漏。")

    # ---------- bottom strips on diagram slides ----------
    refit_pic(slides[I_POS])
    strip(slides[I_POS],
          "一句话定位　① 不是再做一个看板：一条可追溯闭环串到底　② 三支柱 = 数据传输 / 智能体决策 / 可视化操作　"
          "③ 本期边界：模拟 + 虚拟执行器、本地可复现，BearPi 遥测按来源标记接入")
    refit_pic(slides[I_MAIN])
    strip(slides[I_MAIN],
          "硬边界　① 浏览器不直连 MQTT，凭据不下发　② 模拟器不写告警 / 结论　"
          "③ 模型不拼 SQL / MQTT / HTTP，工具白名单 + JSON Schema 校验")
    refit_pic(slides[I_BE])
    strip(slides[I_BE],
          "模块化单体 = 单进程 · 逻辑模块 · 可替换适配器（DataSource / Actuator）；15 天不拆微服务，拆服务是 P2 演进")
    refit_pic(slides[I_FE])
    strip(slides[I_FE],
          "三角色三入口共享一个内核（api.js / live-data.js / roles.js）· 浏览器只走 REST / SSE · "
          "Vue 3 全局运行时 + 模块化原生 JS，未上 TS/SFC 是取舍")

    # 三条主线页：缩图后会悬空的模板小图标/小矩形，直接清掉
    for sh in list(slides[I_MAIN].shapes):
        if sh.name in ("矩形 14", "矩形 15", "矩形 16", "矩形 17") or (
            sh.shape_type == MSO_SHAPE_TYPE.PICTURE and (sh.width or 0) < 100 * PX
        ):
            sh._element.getparent().remove(sh._element)

    # ---------- new slide A: 数据主线九跳 ----------
    layout = slides[I_FE].slide_layout
    sA = prs.slides.add_slide(layout)
    for ph in list(sA.placeholders):
        ph._element.getparent().remove(ph._element)
    clone_chrome(prs, slides[I_FE], sA, "数据主线：一条遥测的九站旅程")

    hops = [
        ("① 产生", ["来源A：SimulationEngine（带种子情景）", "来源B：BearPi E53_IA1 真机适配器",
                    "统一合同：eventId / plotId / 指标 / 质量 / scenarioId"], BLUE2),
        ("② 传输", ["协议 MQTT · 服务端 Mosquitto 2", "客户端 Paho · QoS1 至少送达一次", "真实链路可选，默认进程内直注"], BLUE2),
        ("③ 缓冲", ["Redis Stream agri.telemetry", "消费组 · 批量 50 · DLQ 死信"], BLUE2),
        ("④ 治理", ["校验 · 去重 · 三维质量评分", "新鲜度 / 完整度 / 可信度"], BLUE2),
        ("⑤ 落库", ["PostgreSQL 16 entity_record", "JSONB · Flyway V1–V5 · H2 回退"], BLUE2),
        ("⑥ 计算", ["规则→告警(迟滞/冷却)", "根因→短期预测→就绪度"], TEAL),
        ("⑦ 推送", ["SSE /events/stream", "按用户隔离 · 心跳 · eventId 作 SSE id"], TEAL),
        ("⑧ 呈现", ["Vue 3 工作台秒级刷新", "浏览器不连 MQTT"], TEAL),
        ("⑨ 回执", ["处方→安全门→虚拟执行", "ACK→效果评价→实绩回写"], TEAL),
    ]
    BW, BH = 224, 135
    R1Y, R2Y = 145, 315
    xs1 = [40, 280, 520, 760, 1000]
    for i, (t, b, c) in enumerate(hops[0:5]):
        card(sA, xs1[i], R1Y, BW, BH, t, b, c)
    for gx in (264, 504, 744, 984):
        arrow(sA, gx, R1Y + 57, 16, 22, BLUE2)
    arrow(sA, 1096, 283, 32, 28, TEAL, MSO_SHAPE.DOWN_ARROW)
    xs2 = [1000, 760, 520, 280]
    for i, (t, b, c) in enumerate(hops[5:9]):
        card(sA, xs2[i], R2Y, BW, BH, t, b, c)
    for gx in (984, 744, 504, 264):
        arrow(sA, gx, R2Y + 57, 16, 22, TEAL, rot=180)
    card(sA, 40, R2Y, BW, BH, "↺ 闭环", ["实绩回写 → 计划 / 回放优化", "回到①；固定 scenarioId 可复现"], ORANGE,
         fill=ORANGEBG)
    arrow(sA, 136, 283, 32, 28, ORANGE, MSO_SHAPE.UP_ARROW)

    deg = sA.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(40), px(480), px(1184), px(115))
    deg.adjustments[0] = 0.10
    deg.fill.solid()
    deg.fill.fore_color.rgb = ORANGEBG
    deg.line.color.rgb = ORANGE
    deg.line.width = Pt(1.25)
    deg.shadow.inherit = False
    tf = deg.text_frame
    tf.margin_left = px(14)
    tf.margin_right = px(14)
    tf.margin_top = px(6)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    fill_text(tf, [
        ("降级路径：可见，而不是藏起来（dependencyStatus() 实时上报 UP / DEGRADED）", 12, True, ORANGE),
        ("MQTT 断 → 命令通道 FALLBACK_OR_IDLE，虚拟执行走进程内，页面标注 DEGRADED", 10.5, False, INK),
        ("Redis 断 → 进程内直注 + 内存回退，不谎报事件丢失", 10.5, False, INK),
        ("LLM 断 → rules-only 确定性引擎接管，前端显示降级徽章，不伪装模型结果", 10.5, False, INK),
    ])
    textbox(sA, 40, 606, 1200, 36, [
        ("组件名与代码一一对应：SimulationEngine / hardware/bearpi_e53_bridge.py / MqttCommandGateway / "
         "RedisStreamWorker(agri.telemetry+DLQ) / AgriStore(entity_record) / SseEmitter(/events/stream)，每一站都指得到源码。",
         9, False, GRAY),
    ])
    set_notes(sA,
              "[90s] 这页是本章核心：带评委跟一条遥测走九站。重点展开三站：第三站 Redis Stream 是真用了——消费组 + 批量 50 + "
              "死信队列，失败消息进 DLQ 而不是静默丢；第四站 三维质量评分（新鲜度/完整度/可信度），不达标后面就不给可执行灌溉处方，"
              "监测直接服务“信任”层；第七站 SSE 按用户隔离、带心跳、eventId 作 SSE id，前端去重和回放都靠它。其余扫过。"
              "最后指底部橙框：每一站都有降级路径，且降级可见。时间紧可压到 60s，只展开三四站。")

    # ---------- new slide B: 工程可信度 ----------
    sB = prs.slides.add_slide(layout)
    for ph in list(sB.placeholders):
        ph._element.getparent().remove(ph._element)
    clone_chrome(prs, slides[I_FE], sB, "工程可信度：可复现、可降级、可观测")
    tiles = [
        ("18,413 行", ["后端 Java · 单 Gradle 模块", "“模块化单体”= 逻辑模块、单进程"]),
        ("55,003 行", ["前端 JS+CSS · 34 js / 20 css", "5 个页面：农户/管理员/系统管理员/登录…"]),
        ("157 + 1", ["REST 映射 + SSE 事件流", "统一入口 /api/v1 与 /events/stream"]),
        ("17 文件 · 2,874 行", ["前端测试 + Java JUnit", "双保险，答辩前可跑回归"]),
        ("V1–V5", ["Flyway 数据库迁移", "PG16 生产 / H2 单机一键回退"]),
        ("4 服务", ["docker-compose 一键起", "PG16 + Redis7 + Mosquitto2 + API"]),
        ("2 包 · 8 指标", ["完整 Crop Pack（tomato/cucumber）", "统一指标编码，Schema 校验"]),
        ("3 字段", ["sourceMode / provenance / dataOrigin", "真实与模拟在落库时区分，V5 迁移固化"]),
    ]
    TW, TH = 280, 150
    for i, (num, body) in enumerate(tiles):
        x = 40 + (i % 4) * 304
        y = 150 if i < 4 else 324
        sp = sB.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(TW), px(TH))
        sp.adjustments[0] = 0.08
        sp.fill.solid()
        sp.fill.fore_color.rgb = LIGHT
        sp.line.color.rgb = BLUE
        sp.line.width = Pt(1.25)
        sp.shadow.inherit = False
        tf = sp.text_frame
        tf.margin_left = px(10)
        tf.margin_right = px(10)
        tf.margin_top = px(10)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        fill_text(tf, [(num, 24, True, BLUE)] + [(t, 10, False, RGBColor(0x33, 0x3F, 0x4A)) for t in body])
    strip(sB, "可观测：Actuator + Prometheus 指标 · dependencyStatus() 实时上报 mqtt / redis / 命令通道 UP / DEGRADED",
          y=500, h=64, size=12)
    textbox(sB, 40, 584, 1200, 44, [
        ("“这一页的每个数字都能现场 3 分钟复跑——这就是我们对‘可信’的定义。”", 15, True, ORANGE),
    ], align=PP_ALIGN.CENTER)
    set_notes(sB,
              "[60s] 收尾页，只念数字。后端 1.8 万行 Java、前端 5.5 万行 JS+CSS；157 个 REST 加 1 条 SSE；前端 17 个测试文件加 "
              "Java JUnit；Flyway 五个迁移；compose 一键四服务；两个完整 Crop Pack、八个统一指标编码；三个来源字段由 V5 迁移固化。"
              "收尾句：这一页的每个数字都能现场 3 分钟复跑——这就是我们对“可信”的定义。（超时时本页压到 30s：只念第一排数字和收尾句。）")

    # ---------- appendix note on the dense 6-layer map ----------
    set_notes(slides[I_OLD26],
              "[备用页·不计时间] 评委若追问“能展开讲讲分层吗”翻到这页：按 现场设备 / 边缘接入 / 消息与数据 / 领域服务 / "
              "智能决策 / 应用运维 六层讲，强调每层接口可替换——模拟器可换真设备、H2 可换 PG、rules-only 可换回 LLM。")

    # ---------- reorder ----------
    sldIdLst = prs.slides._sldIdLst
    elems = list(sldIdLst)
    n_new_a = len(elems) - 2  # index of sA
    order = list(range(0, 16)) + [I_POS, I_MAIN, n_new_a, I_BE, I_FE, n_new_a + 1] + list(range(21, 31)) + [I_OLD26]
    assert len(order) == len(elems) == 33, (len(order), len(elems))
    for e in elems:
        sldIdLst.remove(e)
    for i in order:
        sldIdLst.append(elems[i])

    prs.save(DST)
    print("saved", DST)


if __name__ == "__main__":
    main()
